"""اختبار الدورة الكاملة لـ Excel و PowerPoint.

بيتحقق من ادعاءات التصميم الأساسية:
  1. الرسوم البيانية والصور بتفضل موجودة بعد الدمج (ده سبب إننا
     بنعدّل الـ XML بدل ما نستخدم openpyxl للكتابة).
  2. المعادلات مابتتترجمش ومابتتكسرش.
  3. الأرقام مابتتلمسش.
  4. ملاحظات الشرائح بتتلقط (بينساها أغلب الأدوات).
"""
from __future__ import annotations

import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app.tools.translator.formats import pptx_fmt, xlsx_fmt  # noqa: E402
from app.tools.translator.formats.base import (  # noqa: E402
    is_translatable,
    parse_tagged_text,
    strip_tags,
)

SAMPLES = Path("storage/samples")


def fake_translate(text: str) -> str:
    pieces = parse_tagged_text(text)
    out = []
    for tag, chunk in pieces:
        translated = f"[EN:{chunk.strip()}]" if chunk.strip() else chunk
        out.append(f"<g{tag}>{translated}</g{tag}>" if tag else translated)
    return "".join(out)


# ---------------------------------------------------------------------------
def build_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.comments import Comment
    from openpyxl.styles import Alignment, Font

    wb = Workbook()
    ws = wb.active
    ws.title = "الميزانية"
    ws.sheet_view.rightToLeft = True

    rows = [
        ["البند", "الربع الأول", "الربع الثاني"],
        ["إيرادات الاستشارات", 120000, 135000],
        ["تكاليف التشغيل", 45000, 48000],
        ["صافي الربح", None, None],
    ]
    for row in rows:
        ws.append(row)

    # معادلات — لازم تفضل معادلات
    ws["B5"] = "=B2-B3"
    ws["C5"] = "=C2-C3"
    ws["A5"] = "الإجمالي"

    # خلية بتنسيق مختلط
    ws["A7"] = "ملاحظة مهمة للمراجعة"
    ws["A7"].font = Font(bold=True)
    ws["A7"].alignment = Alignment(horizontal="right", readingOrder=2)

    ws["B7"].comment = Comment("راجع الأرقام مع قسم المالية", "المدقق")

    ws2 = wb.create_sheet("البيانات")
    ws2["A1"] = "وصف تفصيلي للبنود المالية."

    # رسم بياني — ده اللي عايزين نتأكد إنه مابيضيعش
    chart = BarChart()
    chart.title = "مقارنة الأرباع"
    data = Reference(ws, min_col=2, min_row=1, max_col=3, max_row=3)
    chart.add_data(data, titles_from_data=True)
    ws.add_chart(chart, "E2")

    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))


def build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "عرض المشروع"
    body = slide.placeholders[1].text_frame
    body.text = "المرحلة الأولى من التنفيذ"
    para = body.add_paragraph()
    run1 = para.add_run()
    run1.text = "الميزانية المعتمدة "
    run2 = para.add_run()
    run2.text = "150 ألف ريال"
    run2.font.bold = True

    # ملاحظات الشريحة
    slide.notes_slide.notes_text_frame.text = "نقطة للنقاش: مدة التنفيذ قابلة للتفاوض."

    # شريحة فيها جدول
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "الجدول الزمني"
    table = slide2.shapes.add_table(
        2, 2, Inches(1), Inches(2), Inches(8), Inches(1.5)
    ).table
    table.cell(0, 0).text = "المرحلة"
    table.cell(0, 1).text = "المدة"
    table.cell(1, 0).text = "التحليل"
    table.cell(1, 1).text = "شهران"

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


# ---------------------------------------------------------------------------
def check_xlsx() -> list[str]:
    failures: list[str] = []
    source = SAMPLES / "budget_ar.xlsx"
    output = SAMPLES / "budget_en.xlsx"
    build_xlsx(source)

    result = xlsx_fmt.extract(source)
    print(f"\n=== Excel: {len(result.units)} وحدة، {result.page_count} ورقة ===")
    for unit in result.units:
        flag = "" if unit.meta.get("translatable", True) else "  [غير مترجَم]"
        print(f"  [{unit.kind:10}] {unit.location:20} | {strip_tags(unit.text)[:40]}{flag}")

    texts = [strip_tags(u.text) for u in result.units]

    # المعادلات مالهاش وحدات
    if any(text.startswith("=") for text in texts):
        failures.append("معادلة اتستخرجت للترجمة — ده بيكسر الملف")

    # الأرقام مالهاش وحدات
    if any(text.strip().replace(",", "").isdigit() for text in texts):
        failures.append("رقم اتستخرج للترجمة")

    # أسماء الأوراق موجودة لكن معلَّمة كغير قابلة للترجمة
    sheet_names = [u for u in result.units if u.kind == "sheet_name"]
    if not sheet_names:
        failures.append("أسماء الأوراق ماتستخرجتش")
    elif any(u.meta.get("translatable", True) for u in sheet_names):
        failures.append("اسم ورقة معلَّم كقابل للترجمة — بيكسر المعادلات")

    if not any(u.kind == "comment" for u in result.units):
        failures.append("تعليق الخلية ماتستخرجش")

    translations = {
        u.unit_key: fake_translate(u.text)
        for u in result.units
        if is_translatable(u.text) and u.meta.get("translatable", True)
    }
    xlsx_fmt.merge(source, output, translations, target_rtl=False)
    print(f"  اتكتب: {output.name} ({output.stat().st_size:,} بايت)، "
          f"{len(translations)} وحدة مترجمة")

    # الادعاء الأساسي: الرسم البياني مابيضيعش
    with zipfile.ZipFile(source) as z:
        charts_before = [n for n in z.namelist() if "charts/chart" in n]
    with zipfile.ZipFile(output) as z:
        names_after = z.namelist()
        charts_after = [n for n in names_after if "charts/chart" in n]
    print(f"  الرسوم البيانية: قبل={len(charts_before)} بعد={len(charts_after)}")
    if len(charts_after) < len(charts_before):
        failures.append("الرسم البياني ضاع بعد الدمج")

    # المعادلات لسه موجودة وسليمة
    from openpyxl import load_workbook
    wb = load_workbook(str(output))
    ws = wb["الميزانية"]
    if ws["B5"].value != "=B2-B3":
        failures.append(f"المعادلة اتكسرت: {ws['B5'].value!r}")
    if ws["B2"].value != 120000:
        failures.append(f"رقم اتغيّر: {ws['B2'].value!r}")
    if not str(ws["A2"].value or "").startswith("[EN:"):
        failures.append(f"النص ماتترجمش: {ws['A2'].value!r}")
    if ws.sheet_view.rightToLeft:
        failures.append("اتجاه الورقة لسه من اليمين لليسار")

    after = xlsx_fmt.extract(output)
    translated = sum(1 for u in after.units if "[EN:" in u.text)
    print(f"  وصلتها الترجمة: {translated}/{len(translations)}")
    if translated < len(translations):
        failures.append(f"وحدات ماوصلتهاش الترجمة: {len(translations) - translated}")

    return failures


def check_pptx() -> list[str]:
    failures: list[str] = []
    source = SAMPLES / "deck_ar.pptx"
    output = SAMPLES / "deck_en.pptx"
    build_pptx(source)

    result = pptx_fmt.extract(source)
    print(f"\n=== PowerPoint: {len(result.units)} وحدة، {result.page_count} شريحة ===")
    for unit in result.units:
        print(f"  [{unit.kind:6}] {unit.location:30} | {strip_tags(unit.text)[:40]}")

    if not any(u.kind == "notes" for u in result.units):
        failures.append("ملاحظات الشريحة ماتستخرجتش")

    if not any("جدول" in u.location for u in result.units):
        failures.append("خلايا الجدول ماتستخرجتش")

    mixed = [u for u in result.units if u.placeholders]
    print(f"  وحدات بتنسيق مختلط: {len(mixed)}")
    if not mixed:
        failures.append("الفقرة متعددة التنسيق ماتوسمتش")

    translations = {
        u.unit_key: fake_translate(u.text)
        for u in result.units
        if is_translatable(u.text)
    }
    pptx_fmt.merge(source, output, translations, target_rtl=False)
    print(f"  اتكتب: {output.name} ({output.stat().st_size:,} بايت)، "
          f"{len(translations)} وحدة مترجمة")

    after = pptx_fmt.extract(output)
    translated = sum(1 for u in after.units if "[EN:" in u.text)
    print(f"  وصلتها الترجمة: {translated}/{len(translations)}")
    if translated < len(translations):
        failures.append(f"وحدات ماوصلتهاش الترجمة: {len(translations) - translated}")

    after_map = {u.unit_key: u for u in after.units}
    for unit in mixed:
        target = after_map.get(unit.unit_key)
        if target and not target.placeholders:
            failures.append(f"ضاع تنسيق: {unit.unit_key}")

    # تقرير الطفح
    report = pptx_fmt.overflow_report(result.units, translations)
    print(f"  أشكال مرشّحة للطفح: {len(report)}")

    return failures


def main() -> int:
    failures = check_xlsx() + check_pptx()

    print("\n" + "=" * 62)
    if failures:
        print(f"فشل: {len(failures)} مشكلة")
        for item in failures:
            print(f"  ✗ {item}")
        return 1
    print("نجح: Excel و PowerPoint سليمين ✓")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
